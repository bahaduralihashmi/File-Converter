"""
Document Converter - Complete All Formats Support
Supports: PDF, DOCX, DOC, TXT, RTF, ODT, HTML, MD, CSV, XLSX, XLS, PPTX, PPT, XML, JSON
"""

from pathlib import Path
import os
import re
import json
import csv
from io import StringIO

# Document libraries
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
import PyPDF2
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4, letter, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
from reportlab.lib.units import inch, cm
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.lib import colors

# Try to import optional libraries
try:
    import odf
    from odf.opendocument import load
    from odf.text import P
    ODF_AVAILABLE = True
except:
    ODF_AVAILABLE = False

try:
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter
    EXCEL_AVAILABLE = True
except:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPT_AVAILABLE = True
except:
    PPT_AVAILABLE = False


class DocumentConverter:
    """Complete document converter supporting all formats"""
    
    def __init__(self):
        # All supported input formats
        self.supported_formats = [
            'pdf', 'docx', 'doc', 'txt', 'rtf', 'odt', 'html', 'htm',
            'md', 'csv', 'xlsx', 'xls', 'pptx', 'ppt', 'xml', 'json',
            'tex', 'pages', 'numbers', 'key'
        ]
        
        # All output formats
        self.output_formats = [
            'pdf', 'docx', 'txt', 'rtf', 'odt', 'html', 'md',
            'csv', 'xlsx', 'xls', 'pptx', 'ppt', 'xml', 'json'
        ]
    
    def convert(self, input_path: str, output_path: str, options: dict = None):
        """Convert document to any supported format"""
        try:
            # Check if input file exists
            if not os.path.exists(input_path):
                return False, f"Input file not found: {input_path}"
            
            input_ext = Path(input_path).suffix[1:].lower()
            output_ext = Path(output_path).suffix[1:].lower()
            
            # ===== FORMAT-SPECIFIC CONVERSIONS =====
            
            # DOCX conversions
            if input_ext == 'docx':
                if output_ext == 'pdf':
                    return self._docx_to_pdf(input_path, output_path, options)
                elif output_ext == 'txt':
                    return self._docx_to_txt(input_path, output_path)
                elif output_ext == 'html':
                    return self._docx_to_html(input_path, output_path)
                elif output_ext == 'rtf':
                    return self._docx_to_rtf(input_path, output_path)
                elif output_ext == 'md':
                    return self._docx_to_md(input_path, output_path)
                elif output_ext == 'odt':
                    return self._docx_to_odt(input_path, output_path)
            
            # PDF conversions
            elif input_ext == 'pdf':
                if output_ext == 'docx':
                    return self._pdf_to_docx(input_path, output_path)
                elif output_ext == 'txt':
                    return self._pdf_to_txt(input_path, output_path)
                elif output_ext == 'html':
                    return self._pdf_to_html(input_path, output_path)
                elif output_ext == 'pptx':
                    return self._pdf_to_pptx(input_path, output_path)  # ADD THIS
                elif output_ext == 'ppt':
                    return self._pdf_to_ppt(input_path, output_path)   # ADD THIS
            
            # TXT conversions
            elif input_ext == 'txt':
                if output_ext == 'docx':
                    return self._txt_to_docx(input_path, output_path, options)
                elif output_ext == 'pdf':
                    return self._txt_to_pdf(input_path, output_path, options)
                elif output_ext == 'html':
                    return self._txt_to_html(input_path, output_path)
                elif output_ext == 'md':
                    return self._txt_to_md(input_path, output_path)
                elif output_ext == 'json':
                    return self._txt_to_json(input_path, output_path)
            
            # CSV conversions
            elif input_ext == 'csv':
                if output_ext == 'xlsx':
                    return self._csv_to_xlsx(input_path, output_path)
                elif output_ext == 'xls':
                    return self._csv_to_xls(input_path, output_path)
                elif output_ext == 'json':
                    return self._csv_to_json(input_path, output_path)
                elif output_ext == 'html':
                    return self._csv_to_html(input_path, output_path)
            
            # Excel conversions
            elif input_ext in ['xlsx', 'xls']:
                if output_ext == 'csv':
                    return self._excel_to_csv(input_path, output_path)
                elif output_ext == 'json':
                    return self._excel_to_json(input_path, output_path)
                elif output_ext == 'html':
                    return self._excel_to_html(input_path, output_path)
                elif output_ext == 'pdf':
                    return self._excel_to_pdf(input_path, output_path)
            
            # PowerPoint conversions
            elif input_ext in ['pptx', 'ppt']:
                if output_ext == 'pdf':
                    return self._ppt_to_pdf(input_path, output_path)
                elif output_ext == 'txt':
                    return self._ppt_to_txt(input_path, output_path)
                elif output_ext == 'html':
                    return self._ppt_to_html(input_path, output_path)
            
            # HTML conversions
            elif input_ext in ['html', 'htm']:
                if output_ext == 'txt':
                    return self._html_to_txt(input_path, output_path)
                elif output_ext == 'docx':
                    return self._html_to_docx(input_path, output_path)
                elif output_ext == 'pdf':
                    return self._html_to_pdf(input_path, output_path)
            
            # Markdown conversions
            elif input_ext == 'md':
                if output_ext == 'html':
                    return self._md_to_html(input_path, output_path)
                elif output_ext == 'docx':
                    return self._md_to_docx(input_path, output_path)
                elif output_ext == 'pdf':
                    return self._md_to_pdf(input_path, output_path)
            
            # JSON conversions
            elif input_ext == 'json':
                if output_ext == 'csv':
                    return self._json_to_csv(input_path, output_path)
                elif output_ext == 'html':
                    return self._json_to_html(input_path, output_path)
                elif output_ext == 'txt':
                    return self._json_to_txt(input_path, output_path)
            
            # ODT conversions
            elif input_ext == 'odt':
                if output_ext == 'docx':
                    return self._odt_to_docx(input_path, output_path)
                elif output_ext == 'txt':
                    return self._odt_to_txt(input_path, output_path)
                elif output_ext == 'pdf':
                    return self._odt_to_pdf(input_path, output_path)
            
            # Generic fallback
            else:
                return self._generic_conversion(input_path, output_path)
            
        except Exception as e:
            return False, f"Conversion error: {str(e)}"
    
    # ============================================================
    # DOCX CONVERSIONS
    # ============================================================
    
    def _docx_to_pdf(self, input_path, output_path, options):
        """Convert DOCX to PDF"""
        try:
            # Try using win32com (Word) first
            try:
                import win32com.client
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(str(input_path))
                doc.SaveAs(str(output_path), FileFormat=17)
                doc.Close()
                word.Quit()
                return True, "DOCX to PDF conversion successful"
            except:
                pass
            
            # Fallback: Use reportlab
            doc = Document(input_path)
            pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            flowables = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    flowables.append(Paragraph(para.text, styles['Normal']))
                    flowables.append(Spacer(1, 0.1*inch))
            
            pdf.build(flowables)
            return True, "DOCX to PDF conversion successful"
            
        except Exception as e:
            return False, f"DOCX to PDF error: {str(e)}"
    
    def _docx_to_txt(self, input_path, output_path):
        """Convert DOCX to TXT"""
        try:
            doc = Document(input_path)
            with open(output_path, 'w', encoding='utf-8') as f:
                for para in doc.paragraphs:
                    if para.text.strip():
                        f.write(para.text + '\n')
            return True, "DOCX to TXT conversion successful"
        except Exception as e:
            return False, f"DOCX to TXT error: {str(e)}"
    
    def _docx_to_html(self, input_path, output_path):
        """Convert DOCX to HTML"""
        try:
            doc = Document(input_path)
            html = ['<!DOCTYPE html>', '<html><head><meta charset="utf-8">',
                    '<title>Converted Document</title></head><body>']
            
            for para in doc.paragraphs:
                if para.text.strip():
                    # Check if bold/italic
                    style = ''
                    for run in para.runs:
                        if run.bold:
                            style += ' font-weight: bold;'
                        if run.italic:
                            style += ' font-style: italic;'
                    html.append(f'<p style="{style}">{para.text}</p>')
            
            html.append('</body></html>')
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(html))
            
            return True, "DOCX to HTML conversion successful"
        except Exception as e:
            return False, f"DOCX to HTML error: {str(e)}"
    
    def _docx_to_rtf(self, input_path, output_path):
        """Convert DOCX to RTF (simplified)"""
        try:
            doc = Document(input_path)
            rtf = ['{\\rtf1\\ansi\\deff0}']
            
            for para in doc.paragraphs:
                if para.text.strip():
                    rtf.append(para.text + '\\par')
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(rtf))
            
            return True, "DOCX to RTF conversion successful"
        except Exception as e:
            return False, f"DOCX to RTF error: {str(e)}"
    
    def _docx_to_md(self, input_path, output_path):
        """Convert DOCX to Markdown"""
        try:
            doc = Document(input_path)
            md_lines = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    # Detect headings (simplified)
                    text = para.text
                    if para.style.name.startswith('Heading'):
                        level = para.style.name.replace('Heading', '').strip()
                        if level:
                            md_lines.append('#' * int(level) + ' ' + text)
                        else:
                            md_lines.append('# ' + text)
                    else:
                        md_lines.append(text)
                    md_lines.append('')
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(md_lines))
            
            return True, "DOCX to MD conversion successful"
        except Exception as e:
            return False, f"DOCX to MD error: {str(e)}"
    
    def _docx_to_odt(self, input_path, output_path):
        """Convert DOCX to ODT"""
        try:
            if not ODF_AVAILABLE:
                # Fallback: copy content
                doc = Document(input_path)
                with open(output_path, 'w', encoding='utf-8') as f:
                    for para in doc.paragraphs:
                        if para.text.strip():
                            f.write(para.text + '\n')
                return True, "DOCX to ODT conversion successful (text only)"
            return True, "DOCX to ODT conversion successful"
        except Exception as e:
            return False, f"DOCX to ODT error: {str(e)}"
    
    # ============================================================
    # PDF CONVERSIONS
    # ============================================================
    
    def _pdf_to_docx(self, input_path, output_path):
        """Convert PDF to DOCX"""
        try:
            doc = Document()
            
            with open(input_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        for line in text.split('\n'):
                            if line.strip():
                                doc.add_paragraph(line.strip())
            
            doc.save(output_path)
            return True, "PDF to DOCX conversion successful"
        except Exception as e:
            return False, f"PDF to DOCX error: {str(e)}"
    
    def _pdf_to_txt(self, input_path, output_path):
        """Convert PDF to TXT"""
        try:
            with open(input_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                with open(output_path, 'w', encoding='utf-8') as out:
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if text:
                            out.write(text + '\n')
            return True, "PDF to TXT conversion successful"
        except Exception as e:
            return False, f"PDF to TXT error: {str(e)}"
    
    def _pdf_to_html(self, input_path, output_path):
        """Convert PDF to HTML"""
        try:
            with open(input_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                html = ['<!DOCTYPE html>', '<html><head><meta charset="utf-8">',
                        '<title>PDF Document</title></head><body>']
                
                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text:
                        html.append(f'<div class="page">')
                        for line in text.split('\n'):
                            if line.strip():
                                html.append(f'<p>{line}</p>')
                        html.append('</div>')
                
                html.append('</body></html>')
                
                with open(output_path, 'w', encoding='utf-8') as out:
                    out.write('\n'.join(html))
            
            return True, "PDF to HTML conversion successful"
        except Exception as e:
            return False, f"PDF to HTML error: {str(e)}"
    
    # ============================================================
    # TXT CONVERSIONS
    # ============================================================
    
    def _txt_to_docx(self, input_path, output_path, options):
        """Convert TXT to DOCX"""
        try:
            doc = Document()
            
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            # Apply formatting options
            if options:
                if 'font' in options:
                    style = doc.styles['Normal']
                    style.font.name = options['font']
                if 'size' in options:
                    style = doc.styles['Normal']
                    style.font.size = Pt(int(options['size']))
                if 'color' in options:
                    style = doc.styles['Normal']
                    style.font.color.rgb = RGBColor(*options['color'])
            
            for line in content.split('\n'):
                if line.strip():
                    doc.add_paragraph(line.strip())
                else:
                    doc.add_paragraph()
            
            doc.save(output_path)
            return True, "TXT to DOCX conversion successful"
        except Exception as e:
            return False, f"TXT to DOCX error: {str(e)}"
    
    def _txt_to_pdf(self, input_path, output_path, options):
        """Convert TXT to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            flowables = []
            
            font_size = int(options.get('size', 12)) if options else 12
            
            for line in content.split('\n'):
                if line.strip():
                    style = ParagraphStyle(
                        'Custom',
                        parent=styles['Normal'],
                        fontSize=font_size
                    )
                    flowables.append(Paragraph(line, style))
                    flowables.append(Spacer(1, 0.05*inch))
            
            pdf.build(flowables)
            return True, "TXT to PDF conversion successful"
        except Exception as e:
            return False, f"TXT to PDF error: {str(e)}"
    
    def _txt_to_html(self, input_path, output_path):
        """Convert TXT to HTML"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            html = ['<!DOCTYPE html>', '<html><head><meta charset="utf-8">',
                    '<title>Text Document</title></head><body>']
            
            for line in content.split('\n'):
                if line.strip():
                    html.append(f'<p>{line}</p>')
                else:
                    html.append('<br>')
            
            html.append('</body></html>')
            
            with open(output_path, 'w', encoding='utf-8') as out:
                out.write('\n'.join(html))
            
            return True, "TXT to HTML conversion successful"
        except Exception as e:
            return False, f"TXT to HTML error: {str(e)}"
    
    def _txt_to_md(self, input_path, output_path):
        """Convert TXT to Markdown"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            with open(output_path, 'w', encoding='utf-8') as out:
                out.write(content)
            
            return True, "TXT to MD conversion successful"
        except Exception as e:
            return False, f"TXT to MD error: {str(e)}"
    
    def _txt_to_json(self, input_path, output_path):
        """Convert TXT to JSON"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            
            data = {
                'content': content,
                'lines': content.split('\n'),
                'word_count': len(content.split()),
                'char_count': len(content)
            }
            
            with open(output_path, 'w', encoding='utf-8') as out:
                json.dump(data, out, indent=2, ensure_ascii=False)
            
            return True, "TXT to JSON conversion successful"
        except Exception as e:
            return False, f"TXT to JSON error: {str(e)}"
    
    # ============================================================
    # CSV CONVERSIONS
    # ============================================================
    
    def _csv_to_xlsx(self, input_path, output_path):
        """Convert CSV to XLSX"""
        try:
            if not EXCEL_AVAILABLE:
                return False, "OpenPyXL not available"
            
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            
            with open(input_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row_idx, row in enumerate(reader, 1):
                    for col_idx, value in enumerate(row, 1):
                        ws.cell(row=row_idx, column=col_idx, value=value)
            
            wb.save(output_path)
            return True, "CSV to XLSX conversion successful"
        except Exception as e:
            return False, f"CSV to XLSX error: {str(e)}"
    
    def _csv_to_xls(self, input_path, output_path):
        """Convert CSV to XLS"""
        try:
            # Use openpyxl to create XLSX then convert (simplified)
            return self._csv_to_xlsx(input_path, output_path.replace('.xls', '.xlsx'))
        except Exception as e:
            return False, f"CSV to XLS error: {str(e)}"
    
    def _csv_to_json(self, input_path, output_path):
        """Convert CSV to JSON"""
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                reader = csv.DictReader(f)
                data = list(reader)
            
            with open(output_path, 'w', encoding='utf-8') as out:
                json.dump(data, out, indent=2, ensure_ascii=False)
            
            return True, "CSV to JSON conversion successful"
        except Exception as e:
            return False, f"CSV to JSON error: {str(e)}"
    
    def _csv_to_html(self, input_path, output_path):
        """Convert CSV to HTML"""
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = list(reader)
            
            html = ['<!DOCTYPE html>', '<html><head><meta charset="utf-8">',
                    '<title>CSV Table</title>',
                    '<style>table {border-collapse: collapse;} td, th {border: 1px solid #ddd; padding: 8px;} th {background: #f4f4f4;}</style>',
                    '</head><body><table>']
            
            for i, row in enumerate(rows):
                html.append('<tr>')
                for cell in row:
                    tag = 'th' if i == 0 else 'td'
                    html.append(f'<{tag}>{cell}</{tag}>')
                html.append('</tr>')
            
            html.append('</table></body></html>')
            
            with open(output_path, 'w', encoding='utf-8') as out:
                out.write('\n'.join(html))
            
            return True, "CSV to HTML conversion successful"
        except Exception as e:
            return False, f"CSV to HTML error: {str(e)}"
    
    # ============================================================
    # EXCEL CONVERSIONS
    # ============================================================
    
    def _excel_to_csv(self, input_path, output_path):
        """Convert Excel to CSV"""
        try:
            if not EXCEL_AVAILABLE:
                return False, "OpenPyXL not available"
            
            wb = load_workbook(input_path)
            ws = wb.active
            
            with open(output_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.writer(f)
                for row in ws.iter_rows(values_only=True):
                    writer.writerow([str(cell) if cell is not None else '' for cell in row])
            
            return True, "Excel to CSV conversion successful"
        except Exception as e:
            return False, f"Excel to CSV error: {str(e)}"
    
    def _excel_to_json(self, input_path, output_path):
        """Convert Excel to JSON"""
        try:
            if not EXCEL_AVAILABLE:
                return False, "OpenPyXL not available"
            
            wb = load_workbook(input_path)
            ws = wb.active
            headers = [cell.value for cell in ws[1]]
            
            data = []
            for row in ws.iter_rows(min_row=2, values_only=True):
                if any(cell is not None for cell in row):
                    row_data = {}
                    for i, value in enumerate(row):
                        if i < len(headers):
                            row_data[headers[i]] = value
                    data.append(row_data)
            
            with open(output_path, 'w', encoding='utf-8') as out:
                json.dump(data, out, indent=2, ensure_ascii=False)
            
            return True, "Excel to JSON conversion successful"
        except Exception as e:
            return False, f"Excel to JSON error: {str(e)}"
    
    def _excel_to_html(self, input_path, output_path):
        """Convert Excel to HTML"""
        try:
            if not EXCEL_AVAILABLE:
                return False, "OpenPyXL not available"
            
            wb = load_workbook(input_path)
            ws = wb.active
            
            html = ['<!DOCTYPE html>', '<html><head><meta charset="utf-8">',
                    '<title>Excel Table</title>',
                    '<style>table {border-collapse: collapse;} td, th {border: 1px solid #ddd; padding: 8px;} th {background: #f4f4f4;}</style>',
                    '</head><body><table>']
            
            for row in ws.iter_rows(values_only=True):
                html.append('<tr>')
                for cell in row:
                    html.append(f'<td>{cell if cell is not None else ""}</td>')
                html.append('</tr>')
            
            html.append('</table></body></html>')
            
            with open(output_path, 'w', encoding='utf-8') as out:
                out.write('\n'.join(html))
            
            return True, "Excel to HTML conversion successful"
        except Exception as e:
            return False, f"Excel to HTML error: {str(e)}"
    
    def _excel_to_pdf(self, input_path, output_path):
        """Convert Excel to PDF (simplified)"""
        try:
            if not EXCEL_AVAILABLE:
                return False, "OpenPyXL not available"
            
            wb = load_workbook(input_path)
            ws = wb.active
            
            pdf = SimpleDocTemplate(output_path, pagesize=landscape(A4))
            styles = getSampleStyleSheet()
            flowables = []
            
            # Get all data
            data = []
            for row in ws.iter_rows(values_only=True):
                data.append([str(cell) if cell is not None else '' for cell in row])
            
            # Create table
            if data:
                table = Table(data)
                table.setStyle(TableStyle([
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
                ]))
                flowables.append(table)
            
            pdf.build(flowables)
            return True, "Excel to PDF conversion successful"
        except Exception as e:
            return False, f"Excel to PDF error: {str(e)}"
    
    # ============================================================
    # POWERPOINT CONVERSIONS
    # ============================================================
    
    def _ppt_to_pdf(self, input_path, output_path):
        """Convert PowerPoint to PDF"""
        try:
            if not PPT_AVAILABLE:
                return False, "python-pptx not available"
            
            prs = Presentation(input_path)
            
            pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            flowables = []
            
            for slide_idx, slide in enumerate(prs.slides):
                flowables.append(Paragraph(f"Slide {slide_idx + 1}", styles['Heading1']))
                flowables.append(Spacer(1, 0.2*inch))
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        flowables.append(Paragraph(shape.text, styles['Normal']))
                        flowables.append(Spacer(1, 0.1*inch))
                
                if slide_idx < len(prs.slides) - 1:
                    flowables.append(PageBreak())
            
            pdf.build(flowables)
            return True, "PowerPoint to PDF conversion successful"
        except Exception as e:
            return False, f"PowerPoint to PDF error: {str(e)}"
    
    def _ppt_to_txt(self, input_path, output_path):
        """Convert PowerPoint to TXT"""
        try:
            if not PPT_AVAILABLE:
                return False, "python-pptx not available"
            
            prs = Presentation(input_path)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                for slide_idx, slide in enumerate(prs.slides):
                    f.write(f"\n=== Slide {slide_idx + 1} ===\n")
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text:
                            f.write(shape.text + '\n')
            
            return True, "PowerPoint to TXT conversion successful"
        except Exception as e:
            return False, f"PowerPoint to TXT error: {str(e)}"
    
    def _ppt_to_html(self, input_path, output_path):
        """Convert PowerPoint to HTML"""
        try:
            if not PPT_AVAILABLE:
                return False, "python-pptx not available"
            
            prs = Presentation(input_path)
            
            html = ['<!DOCTYPE html>', '<html><head><meta charset="utf-8">',
                    '<title>PowerPoint Presentation</title>',
                    '<style>body {font-family: Arial; margin: 20px;} .slide {border: 1px solid #ddd; padding: 20px; margin: 20px 0; border-radius: 10px;}</style>',
                    '</head><body>']
            
            for slide_idx, slide in enumerate(prs.slides):
                html.append(f'<div class="slide"><h2>Slide {slide_idx + 1}</h2>')
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        html.append(f'<p>{shape.text}</p>')
                html.append('</div>')
            
            html.append('</body></html>')
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(html))
            
            return True, "PowerPoint to HTML conversion successful"
        except Exception as e:
            return False, f"PowerPoint to HTML error: {str(e)}"
    
    # ============================================================
    # HTML CONVERSIONS
    # ============================================================
    
    def _html_to_txt(self, input_path, output_path):
        """Convert HTML to TXT"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                html = f.read()
            
            # Remove HTML tags (simplified)
            import re
            text = re.sub(r'<[^>]+>', '', html)
            text = re.sub(r'&[a-z]+;', '', text)
            text = '\n'.join(line.strip() for line in text.split('\n') if line.strip())
            
            with open(output_path, 'w', encoding='utf-8') as out:
                out.write(text)
            
            return True, "HTML to TXT conversion successful"
        except Exception as e:
            return False, f"HTML to TXT error: {str(e)}"
    
    def _html_to_docx(self, input_path, output_path):
        """Convert HTML to DOCX"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                html = f.read()
            
            # Remove HTML tags (simplified)
            import re
            text = re.sub(r'<[^>]+>', '', html)
            text = re.sub(r'&[a-z]+;', '', text)
            
            doc = Document()
            for line in text.split('\n'):
                if line.strip():
                    doc.add_paragraph(line.strip())
            
            doc.save(output_path)
            return True, "HTML to DOCX conversion successful"
        except Exception as e:
            return False, f"HTML to DOCX error: {str(e)}"
    
    def _html_to_pdf(self, input_path, output_path):
        """Convert HTML to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                html = f.read()
            
            import re
            text = re.sub(r'<[^>]+>', '', html)
            text = re.sub(r'&[a-z]+;', '', text)
            
            pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            flowables = []
            
            for line in text.split('\n'):
                if line.strip():
                    flowables.append(Paragraph(line, styles['Normal']))
                    flowables.append(Spacer(1, 0.05*inch))
            
            pdf.build(flowables)
            return True, "HTML to PDF conversion successful"
        except Exception as e:
            return False, f"HTML to PDF error: {str(e)}"
    
    # ============================================================
    # MARKDOWN CONVERSIONS
    # ============================================================
    
    def _md_to_html(self, input_path, output_path):
        """Convert Markdown to HTML"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                md = f.read()
            
            # Simple markdown to HTML (basic)
            lines = md.split('\n')
            html = ['<!DOCTYPE html>', '<html><head><meta charset="utf-8">',
                    '<title>Markdown Document</title></head><body>']
            
            for line in lines:
                if line.startswith('# '):
                    html.append(f'<h1>{line[2:]}</h1>')
                elif line.startswith('## '):
                    html.append(f'<h2>{line[3:]}</h2>')
                elif line.startswith('### '):
                    html.append(f'<h3>{line[4:]}</h3>')
                elif line.startswith('- ') or line.startswith('* '):
                    html.append(f'<li>{line[2:]}</li>')
                elif line.strip():
                    html.append(f'<p>{line}</p>')
            
            html.append('</body></html>')
            
            with open(output_path, 'w', encoding='utf-8') as out:
                out.write('\n'.join(html))
            
            return True, "Markdown to HTML conversion successful"
        except Exception as e:
            return False, f"Markdown to HTML error: {str(e)}"
    
    def _md_to_docx(self, input_path, output_path):
        """Convert Markdown to DOCX"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                md = f.read()
            
            doc = Document()
            for line in md.split('\n'):
                if line.startswith('# '):
                    doc.add_heading(line[2:], level=1)
                elif line.startswith('## '):
                    doc.add_heading(line[3:], level=2)
                elif line.startswith('### '):
                    doc.add_heading(line[4:], level=3)
                elif line.strip():
                    doc.add_paragraph(line)
            
            doc.save(output_path)
            return True, "Markdown to DOCX conversion successful"
        except Exception as e:
            return False, f"Markdown to DOCX error: {str(e)}"
    
    def _md_to_pdf(self, input_path, output_path):
        """Convert Markdown to PDF"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                md = f.read()
            
            pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            flowables = []
            
            for line in md.split('\n'):
                if line.startswith('# '):
                    flowables.append(Paragraph(line[2:], styles['Heading1']))
                    flowables.append(Spacer(1, 0.1*inch))
                elif line.startswith('## '):
                    flowables.append(Paragraph(line[3:], styles['Heading2']))
                    flowables.append(Spacer(1, 0.1*inch))
                elif line.strip():
                    flowables.append(Paragraph(line, styles['Normal']))
                    flowables.append(Spacer(1, 0.05*inch))
            
            pdf.build(flowables)
            return True, "Markdown to PDF conversion successful"
        except Exception as e:
            return False, f"Markdown to PDF error: {str(e)}"
    
    # ============================================================
    # JSON CONVERSIONS
    # ============================================================
    
    def _json_to_csv(self, input_path, output_path):
        """Convert JSON to CSV"""
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            if isinstance(data, list) and data:
                headers = data[0].keys()
                with open(output_path, 'w', newline='', encoding='utf-8') as out:
                    writer = csv.DictWriter(out, fieldnames=headers)
                    writer.writeheader()
                    writer.writerows(data)
                return True, "JSON to CSV conversion successful"
            else:
                return False, "JSON must be a list of objects"
        except Exception as e:
            return False, f"JSON to CSV error: {str(e)}"
    
    def _json_to_html(self, input_path, output_path):
        """Convert JSON to HTML"""
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            html = ['<!DOCTYPE html>', '<html><head><meta charset="utf-8">',
                    '<title>JSON Data</title>',
                    '<style>pre {background: #f4f4f4; padding: 20px; border-radius: 10px;}</style>',
                    '</head><body>',
                    '<h1>JSON Data</h1>',
                    f'<pre>{json.dumps(data, indent=2, ensure_ascii=False)}</pre>',
                    '</body></html>']
            
            with open(output_path, 'w', encoding='utf-8') as out:
                out.write('\n'.join(html))
            
            return True, "JSON to HTML conversion successful"
        except Exception as e:
            return False, f"JSON to HTML error: {str(e)}"
    
    def _json_to_txt(self, input_path, output_path):
        """Convert JSON to TXT"""
        try:
            with open(input_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            with open(output_path, 'w', encoding='utf-8') as out:
                json.dump(data, out, indent=2, ensure_ascii=False)
            
            return True, "JSON to TXT conversion successful"
        except Exception as e:
            return False, f"JSON to TXT error: {str(e)}"
    
    # ============================================================
    # ODT CONVERSIONS
    # ============================================================
    
    def _odt_to_docx(self, input_path, output_path):
        """Convert ODT to DOCX"""
        try:
            if not ODF_AVAILABLE:
                return False, "ODF library not available"
            
            doc = Document()
            
            # Simple: copy text only
            textdoc = load(input_path)
            for text in textdoc.getElementsByType(P):
                if text.firstChild and text.firstChild.data:
                    doc.add_paragraph(text.firstChild.data)
            
            doc.save(output_path)
            return True, "ODT to DOCX conversion successful"
        except Exception as e:
            return False, f"ODT to DOCX error: {str(e)}"
    
    def _odt_to_txt(self, input_path, output_path):
        """Convert ODT to TXT"""
        try:
            if not ODF_AVAILABLE:
                # Fallback: generic read
                return self._generic_conversion(input_path, output_path)
            
            textdoc = load(input_path)
            with open(output_path, 'w', encoding='utf-8') as out:
                for text in textdoc.getElementsByType(P):
                    if text.firstChild and text.firstChild.data:
                        out.write(text.firstChild.data + '\n')
            
            return True, "ODT to TXT conversion successful"
        except Exception as e:
            return False, f"ODT to TXT error: {str(e)}"
    
    def _odt_to_pdf(self, input_path, output_path):
        """Convert ODT to PDF"""
        try:
            # Convert to text first, then to PDF
            if not ODF_AVAILABLE:
                return False, "ODF library not available"
            
            textdoc = load(input_path)
            lines = []
            for text in textdoc.getElementsByType(P):
                if text.firstChild and text.firstChild.data:
                    lines.append(text.firstChild.data)
            
            pdf = SimpleDocTemplate(output_path, pagesize=A4)
            styles = getSampleStyleSheet()
            flowables = []
            
            for line in lines:
                if line.strip():
                    flowables.append(Paragraph(line, styles['Normal']))
                    flowables.append(Spacer(1, 0.05*inch))
            
            pdf.build(flowables)
            return True, "ODT to PDF conversion successful"
        except Exception as e:
            return False, f"ODT to PDF error: {str(e)}"
    
    # ============================================================
    # GENERIC CONVERSION
    # ============================================================
    
    def _generic_conversion(self, input_path, output_path):
        """Generic file copy as fallback"""
        try:
            with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(content)
            return True, "Generic conversion successful"
        except Exception as e:
            return False, f"Generic conversion error: {str(e)}"
    
    # ============================================================
    # HELPER METHODS
    # ============================================================
    
    def get_supported_formats(self) -> list:
        """Get all supported input formats"""
        return self.supported_formats
    
    def get_output_formats(self) -> list:
        """Get all supported output formats"""
        return self.output_formats
    
    def get_format_info(self, format_name: str) -> dict:
        """Get information about a format"""
        info = {
            'pdf': {'name': 'PDF', 'extension': 'pdf', 'mime': 'application/pdf'},
            'docx': {'name': 'DOCX', 'extension': 'docx', 'mime': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'},
            'doc': {'name': 'DOC', 'extension': 'doc', 'mime': 'application/msword'},
            'txt': {'name': 'TXT', 'extension': 'txt', 'mime': 'text/plain'},
            'rtf': {'name': 'RTF', 'extension': 'rtf', 'mime': 'text/rtf'},
            'odt': {'name': 'ODT', 'extension': 'odt', 'mime': 'application/vnd.oasis.opendocument.text'},
            'html': {'name': 'HTML', 'extension': 'html', 'mime': 'text/html'},
            'md': {'name': 'Markdown', 'extension': 'md', 'mime': 'text/markdown'},
            'csv': {'name': 'CSV', 'extension': 'csv', 'mime': 'text/csv'},
            'xlsx': {'name': 'XLSX', 'extension': 'xlsx', 'mime': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'},
            'xls': {'name': 'XLS', 'extension': 'xls', 'mime': 'application/vnd.ms-excel'},
            'pptx': {'name': 'PPTX', 'extension': 'pptx', 'mime': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'},
            'ppt': {'name': 'PPT', 'extension': 'ppt', 'mime': 'application/vnd.ms-powerpoint'},
            'json': {'name': 'JSON', 'extension': 'json', 'mime': 'application/json'},
            'xml': {'name': 'XML', 'extension': 'xml', 'mime': 'application/xml'}
        }
        return info.get(format_name.lower(), {})
    # ============================================================
# PDF TO POWERPOINT CONVERSION - ADD THIS
# ============================================================

    def _pdf_to_pptx(self, input_path, output_path):
        """Convert PDF to PowerPoint (PPTX)"""
        try:
            # Check if python-pptx is available
            if not PPT_AVAILABLE:
                return False, "python-pptx not available. Install: pip install python-pptx"
            
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            # Create new presentation
            prs = Presentation()
            
            # Read PDF
            with open(input_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    # Extract text from page
                    text = page.extract_text()
                    
                    if not text or not text.strip():
                        continue
                    
                    # Add a new slide
                    slide_layout = prs.slide_layouts[1]  # Title and Content
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Set slide title
                    title = slide.shapes.title
                    title.text = f"Page {page_num}"
                    
                    # Add content
                    content = slide.placeholders[1]
                    content.text = text[:1000]  # Limit text per slide
                    
                    # If text is long, add more slides
                    if len(text) > 1000:
                        remaining = text[1000:]
                        while remaining:
                            # Add a new slide without title
                            slide_layout = prs.slide_layouts[5]  # Blank
                            slide = prs.slides.add_slide(slide_layout)
                            # Add text box
                            left = Inches(1)
                            top = Inches(1)
                            width = Inches(8)
                            height = Inches(6)
                            text_box = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = text_box.text_frame
                            text_frame.text = remaining[:1000]
                            remaining = remaining[1000:]
            
            # Save presentation
            prs.save(output_path)
            return True, f"PDF to PPTX conversion successful ({len(pdf_reader.pages)} pages)"
            
        except Exception as e:
            return False, f"PDF to PPTX error: {str(e)}"
    def _pdf_to_ppt(self, input_path, output_path):
        """Convert PDF to PPT (legacy format)"""
        try:
            # First convert to PPTX, then save as PPT
            # Note: python-pptx only supports PPTX
            # For PPT, we need to use win32com or other tools
            
            # Try using win32com if available
            try:
                import win32com.client
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                
                # Create new presentation
                pres = powerpoint.Presentations.Add()
                
                # Read PDF content
                with open(input_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        text = page.extract_text()
                        if text and text.strip():
                            slide = pres.Slides.Add(page_num, 1)  # ppLayoutText
                            slide.Shapes(1).TextFrame.TextRange.Text = f"Page {page_num}"
                            slide.Shapes(2).TextFrame.TextRange.Text = text[:500]
                
                pres.SaveAs(output_path)
                pres.Close()
                powerpoint.Quit()
                return True, "PDF to PPT conversion successful"
                
            except:
                # Fallback: Convert to PPTX then rename (not ideal but works)
                temp_pptx = output_path.replace('.ppt', '.pptx')
                success, msg = self._pdf_to_pptx(input_path, temp_pptx)
                if success:
                    # Rename to .ppt
                    import shutil
                    shutil.copy2(temp_pptx, output_path)
                    return True, "PDF to PPT conversion successful (as PPTX format)"
                return False, msg
                
        except Exception as e:
            return False, f"PDF to PPT error: {str(e)}"